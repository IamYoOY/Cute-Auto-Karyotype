
import tkinter as tk
import ctypes
from tkinter import filedialog, messagebox
from pathlib import Path
import os, csv, copy, math
# DPI fix for Windows display scaling: helps Tkinter mouse coordinates match the canvas.
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except Exception:
    try:
        ctypes.windll.user32.SetProcessDPIAware()
    except Exception:
        pass

import cv2
import numpy as np
from PIL import Image, ImageTk, ImageDraw, ImageFont

try:
    from tkinterdnd2 import DND_FILES, TkinterDnD
    DND_OK = True
except Exception:
    DND_OK = False

APP_TITLE = "Cute Auto Karyotype✨"

def imread_unicode(path):
    data = np.fromfile(str(path), dtype=np.uint8)
    return cv2.imdecode(data, cv2.IMREAD_COLOR)

def rotate_bound_white(image, angle):
    h, w = image.shape[:2]
    center = (w/2, h/2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    cos, sin = abs(M[0,0]), abs(M[0,1])
    nw = int(h*sin + w*cos)
    nh = int(h*cos + w*sin)
    M[0,2] += nw/2 - center[0]
    M[1,2] += nh/2 - center[1]
    return cv2.warpAffine(image, M, (nw, nh), flags=cv2.INTER_CUBIC, borderValue=(255,255,255))

def trim_white(img, margin=18):
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    mask = gray < 248
    ys, xs = np.where(mask)
    if len(xs) == 0:
        return img
    x1 = max(xs.min()-margin, 0)
    y1 = max(ys.min()-margin, 0)
    x2 = min(xs.max()+margin, img.shape[1]-1)
    y2 = min(ys.max()+margin, img.shape[0]-1)
    return img[y1:y2+1, x1:x2+1]

def detect_candidates(img, max_count=40):
    h, w = img.shape[:2]
    hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV)
    m1 = cv2.inRange(hsv, np.array([125,18,35]), np.array([179,255,255]))
    m2 = cv2.inRange(hsv, np.array([0,18,35]), np.array([15,255,255]))
    mask = cv2.bitwise_or(m1, m2)
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, np.ones((3,3), np.uint8), iterations=1)
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, cv2.getStructuringElement(cv2.MORPH_ELLIPSE,(15,15)), iterations=2)

    num, labels, stats, centroids = cv2.connectedComponentsWithStats(mask, 8)
    candidates = []
    image_area = h*w
    for i in range(1, num):
        x,y,bw,bh,area = stats[i]
        cx,cy = centroids[i]
        if area < 200: continue
        if area > image_area*0.05: continue
        if (x < 3 or y < 3 or x+bw > w-3 or y+bh > h-3) and area > 2500: continue
        length = max(bw,bh)
        if length < 20: continue
        aspect = max(bw,bh) / max(1,min(bw,bh))

        # V17.1 noise guard:
        # Round WBC/RBC/stain blobs are often large, nearly circular, and should not become chromosome #1.
        # Chromosomes are usually elongated; even small chromosomes rarely look like a huge round disk.
        blob_density = area / max(1, bw*bh)
        is_big_round_blob = (area > 9000 and aspect < 1.65 and blob_density > 0.35)
        is_huge_stain_cloud = (area > 18000 and aspect < 2.0)
        if is_big_round_blob or is_huge_stain_cloud:
            continue

        score = area + 50*length + 300*aspect
        candidates.append(dict(id=0, x=int(x), y=int(y), w=int(bw), h=int(bh),
            cx=float(cx), cy=float(cy), area=int(area), length=int(length), score=float(score),
            selected=True, manual=False, angle=0.0))
    candidates = sorted(candidates, key=lambda d: d["score"], reverse=True)[:max_count]
    candidates = sorted(candidates, key=lambda d: d["length"], reverse=True)
    renumber(candidates)
    return candidates

def renumber(candidates):
    # selected first; ids match crop panel order
    selected = [c for c in candidates if c.get("selected", True)]
    rejected = [c for c in candidates if not c.get("selected", True)]
    candidates[:] = selected + rejected
    for i,c in enumerate(candidates, start=1):
        c["id"] = i
        c["cx"] = c["x"] + c["w"]/2
        c["cy"] = c["y"] + c["h"]/2
        c["length"] = max(c["w"], c["h"])
        c.setdefault("angle", 0.0)
        c.setdefault("manual", False)

def assign_by_length(candidates):
    selected = sorted([c for c in candidates if c["selected"]], key=lambda c: c["length"], reverse=True)
    rejected = [c for c in candidates if not c["selected"]]
    candidates[:] = selected + rejected
    renumber(candidates)

def assign_by_position(candidates):
    selected = sorted([c for c in candidates if c["selected"]], key=lambda c: (round(c["cy"]/80), c["cx"]))
    rejected = [c for c in candidates if not c["selected"]]
    candidates[:] = selected + rejected
    renumber(candidates)

def selected_sorted(candidates):
    return sorted([c for c in candidates if c["selected"]], key=lambda c: c["id"])

def crop_from_original(original, c, pad=40):
    x1 = max(int(c["x"])-pad, 0)
    y1 = max(int(c["y"])-pad, 0)
    x2 = min(int(c["x"]+c["w"])+pad, original.shape[1])
    y2 = min(int(c["y"]+c["h"])+pad, original.shape[0])
    crop = original[y1:y2, x1:x2].copy()
    safe = cv2.copyMakeBorder(crop, 60, 60, 60, 60, cv2.BORDER_CONSTANT, value=(255,255,255))
    rotated = rotate_bound_white(safe, c.get("angle",0.0))
    return trim_white(rotated, margin=20)


def get_label_font(size=34):
    # Try common Windows fonts first. Fallback to PIL default.
    for font_path in [
        "C:/Windows/Fonts/arialbd.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibrib.ttf",
        "C:/Windows/Fonts/calibri.ttf",
    ]:
        try:
            return ImageFont.truetype(font_path, size=size)
        except Exception:
            pass
    try:
        return ImageFont.load_default()
    except Exception:
        return None

def draw_overlay(img, candidates, show_numbers=True):
    pil = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
    draw = ImageDraw.Draw(pil)
    font = get_label_font(72)  # BIG original label font; change 72 to 84/96 if needed

    for c in candidates:
        x,y,w,h = c["x"], c["y"], c["w"], c["h"]
        if c["selected"] and c.get("manual", False): 
            color = (30,90,255)
        elif c["selected"]: 
            color = (0,170,0)
        else: 
            color = (220,30,30)

        # thicker chromosome box
        draw.rectangle([x,y,x+w,y+h], outline=color, width=9)

        if show_numbers:
            label = str(c["id"])

            # measure text
            try:
                bbox = draw.textbbox((0,0), label, font=font)
                tw, th = bbox[2]-bbox[0], bbox[3]-bbox[1]
            except Exception:
                tw, th = 28*len(label), 36

            pad_x, pad_y = 14, 9
            box_w = tw + pad_x*2
            box_h = th + pad_y*2

            # Put label above box, or inside if too close to top
            lx = max(0, x - 8)
            ly = y - box_h - 10
            if ly < 0:
                ly = y + 6

            # shadow/white outline behind label box
            draw.rounded_rectangle(
                [lx-3, ly-3, lx+box_w+3, ly+box_h+3],
                radius=10,
                fill=(255,255,255),
                outline=(255,255,255),
                width=3
            )
            draw.rounded_rectangle(
                [lx, ly, lx+box_w, ly+box_h],
                radius=10,
                fill=color,
                outline=(255,255,255),
                width=3
            )
            draw.text((lx+pad_x, ly+pad_y-3), label, fill=(255,255,255), font=font)
    return pil


def crop_to_transparent_rgba(crop_bgr, threshold=248):
    """Convert white background crop to transparent PNG layer."""
    rgb = cv2.cvtColor(crop_bgr, cv2.COLOR_BGR2RGB)
    gray = cv2.cvtColor(crop_bgr, cv2.COLOR_BGR2GRAY)
    alpha = np.where(gray < threshold, 255, 0).astype(np.uint8)
    # Smooth edge a little, but keep background transparent
    alpha = cv2.medianBlur(alpha, 3)
    rgba = np.dstack([rgb, alpha])
    return Image.fromarray(rgba, mode="RGBA")


def make_svg_layer_file(svg_path, crop_items, cell_w, cell_h, cols):
    """Create an SVG that keeps each chromosome as a separate image object/layer-like element."""
    import base64, io
    rows = math.ceil(len(crop_items)/cols)
    width, height = cols*cell_w, rows*cell_h + 95
    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
             '<rect width="100%" height="100%" fill="white"/>',
             f'<text x="30" y="45" font-family="Arial" font-size="20">Cute Auto Karyotype V17.2 | editable SVG layers | selected: {len(crop_items)}</text>']
    for i, (cid, crop) in enumerate(crop_items):
        r, col = divmod(i, cols)
        x0, y0 = col*cell_w, 85+r*cell_h
        rgba = crop_to_transparent_rgba(crop)
        bio = io.BytesIO(); rgba.save(bio, format="PNG")
        b64 = base64.b64encode(bio.getvalue()).decode("ascii")
        x = x0 + (cell_w - rgba.width)//2
        y = y0 + 25
        parts.append(f'<g id="chromosome_{cid:02d}">')
        parts.append(f'<rect x="{x0+10}" y="{y0+10}" width="{cell_w-20}" height="{cell_h-20}" rx="18" fill="#fcfcfc" stroke="#dddddd" stroke-width="2"/>')
        parts.append(f'<image x="{x}" y="{y}" width="{rgba.width}" height="{rgba.height}" href="data:image/png;base64,{b64}"/>')
        parts.append(f'<text x="{x0+cell_w//2}" y="{y0+cell_h-25}" text-anchor="middle" font-family="Arial" font-size="18">{cid}</text>')
        parts.append('</g>')
    parts.append('</svg>')
    svg_path.write_text("\n".join(parts), encoding="utf-8")

class App:
    def __init__(self, root):
        self.root = root
        self.root.title(APP_TITLE)
        # V17: safer DPI handling for Windows display scaling
        try:
            import ctypes
            ctypes.windll.shcore.SetProcessDpiAwareness(1)
        except Exception:
            pass
        self.root.geometry("1700x980")
        self.root.configure(bg="#fff7fb")

        self.image_path = None
        self.original = None
        self.review_img = None
        self.candidates = []
        self.crop_edits = {}
        self.crop_masks = {}  # persistent white erase mask per chromosome ID
        self.active_id = None
        self.mode = tk.StringVar(value="toggle")

        self.display_scale = 1.0
        self.crop_display_scale = 1.0
        self.left_zoom = tk.IntVar(value=50)
        self.right_zoom = tk.IntVar(value=50)

        self.drag_start = None
        self.temp_rect = None
        self.crop_lasso = []
        self.crop_lasso_item = None
        self.crop_circle_item = None
        self.erase_snapshot_taken = False

        self.undo_stack = []
        self.redo_stack = []
        self.brush_size = tk.IntVar(value=28)
        self.crop_brush_size = tk.IntVar(value=24)
        self.crop_angle_var = tk.StringVar(value="0")
        self.show_numbers = tk.BooleanVar(value=True)
        self.output_dir = Path.cwd() / "karyotype_output"

        self.thumb_refs = []
        self.crop_editor_ref = None

        # V17 final karyotype board state
        self.board_win = None
        self.board_canvas = None
        self.board_items = []
        self.board_refs = []
        self.board_active = None
        self.board_drag = None

        self.build_ui()
        self.set_toggle_mode()

    def build_ui(self):
        tk.Label(self.root, text=APP_TITLE, font=("Arial",24,"bold"), bg="#fff7fb").pack(pady=4)
        tk.Label(self.root, text="Numbers on the original image match the Crop Field IDs. Use Auto Number to synchronize the order.", font=("Arial",11), bg="#fff7fb").pack()

        content = tk.Frame(self.root, bg="#fff7fb")
        content.pack(fill="both", expand=True, padx=6, pady=4)

        # LEFT SIDEBAR
        side = tk.Frame(content, bg="#fff7fb", width=210)
        side.pack(side="left", fill="y", padx=(0,6))

        self.dropbox = tk.Label(side, text="Loaded image\n\nClick / Drop image", width=22, height=5, bg="white", relief="ridge", bd=2, font=("Arial",11))
        self.dropbox.pack(pady=5)
        self.dropbox.bind("<Button-1>", self.choose)
        if DND_OK:
            self.dropbox.drop_target_register(DND_FILES)
            self.dropbox.dnd_bind("<<Drop>>", self.drop)

        info = tk.LabelFrame(side, text="Image info", bg="#fff7fb")
        info.pack(fill="x", pady=5)
        self.info_label = tk.Label(info, text="No image", bg="#fff7fb", justify="left")
        self.info_label.pack(anchor="w", padx=6, pady=4)

        tk.Label(side, text="Max candidates", bg="#fff7fb").pack(anchor="w")
        self.max_var = tk.StringVar(value="100")
        tk.Entry(side, textvariable=self.max_var, width=8).pack(anchor="w")
        self.btn(side,"Detect","#d8ffd8",self.detect).pack(fill="x", pady=2)
        # Select top N is hidden. Use Keep/Reject and Delete rejected instead.
        self.n_var = tk.StringVar(value="")

        tools = tk.LabelFrame(side, text="Tools", bg="#fff7fb")
        tools.pack(fill="x", pady=6)
        self.toggle_btn = self.btn(tools,"Select / Keep-Reject ✅","#d9ffd9",self.set_toggle_mode); self.toggle_btn.pack(fill="x", pady=1)
        self.add_btn = self.btn(tools,"Add Box ✂️",None,self.set_add_mode); self.add_btn.pack(fill="x", pady=1)
        self.erase_btn = self.btn(tools,"Erase Original 🧽",None,self.set_erase_original_mode); self.erase_btn.pack(fill="x", pady=1)
        tk.Label(tools, text="Orig brush", bg="#fff7fb").pack(anchor="w")
        tk.Entry(tools, textvariable=self.brush_size, width=8).pack(anchor="w")

        self.btn(side,"Undo ↶","#fff2cc",self.undo).pack(fill="x", pady=2)
        self.btn(side,"Redo ↷","#fff2cc",self.redo).pack(fill="x", pady=2)
        self.btn(side,"Delete rejected",None,self.delete_rejected).pack(fill="x", pady=2)

        assign = tk.LabelFrame(side, text="Assign Number", bg="#fff7fb")
        assign.pack(fill="x", pady=6)
        self.btn(assign,"Auto Number (Chromosome Size)",None,self.assign_length).pack(fill="x", pady=1)
        self.btn(assign,"Auto Number (Position)",None,self.assign_position).pack(fill="x", pady=1)
        tk.Checkbutton(assign, text="SHOW BIG NUMBERS on original", variable=self.show_numbers, bg="#fff7fb", command=self.refresh_all).pack(anchor="w")

        export = tk.LabelFrame(side, text="Export", bg="#fff7fb")
        export.pack(fill="x", pady=6)
        self.btn(export,"Create Final Board 🧬","#d9f0ff",self.create_final_board).pack(fill="x", pady=2)
        self.btn(export,"Export PNG/PDF 💖","#c9f7ff",self.export_all).pack(fill="x", pady=2)
        self.btn(export,"Open Output 📁",None,self.open_folder).pack(fill="x", pady=2)
        self.btn(export,"Help / Manual PDF", "#e8f4ff", self.open_help_window).pack(fill="x", pady=2)
        self.btn(export,"About", "#f2f2f2", self.open_about_window).pack(fill="x", pady=2)

        # CENTER ORIGINAL
        center = tk.Frame(content, bg="#fff7fb")
        center.pack(side="left", fill="both", expand=True)

        zoombar = tk.Frame(center, bg="#fff7fb")
        zoombar.pack(fill="x")
        tk.Label(zoombar, text="Original image - assign number field", bg="#fff7fb", font=("Arial",11,"bold")).pack(side="left")
        tk.Label(zoombar, text="  Zoom 0–100%", bg="#fff7fb").pack(side="left")
        tk.Scale(zoombar, from_=0, to=100, orient="horizontal", variable=self.left_zoom, command=lambda e:self.refresh_canvas(), length=180).pack(side="left")
        tk.Button(zoombar, text="Fit", command=lambda:self.set_left_zoom(50)).pack(side="left", padx=3)
        tk.Button(zoombar, text="100%", command=lambda:self.set_left_zoom(100)).pack(side="left", padx=3)

        self.status = tk.Label(center, text="Ready", bg="#fff7fb", font=("Arial",10))
        self.status.pack(fill="x", pady=2)

        canvas_wrap = tk.Frame(center, bg="#fff7fb")
        canvas_wrap.pack(fill="both", expand=True)
        self.canvas = tk.Canvas(canvas_wrap, width=850, height=720, bg="white", highlightthickness=1, highlightbackground="#bbbbbb")
        self.canvas_x = tk.Scrollbar(canvas_wrap, orient="horizontal", command=self.canvas.xview)
        self.canvas_y = tk.Scrollbar(canvas_wrap, orient="vertical", command=self.canvas.yview)
        self.canvas.configure(xscrollcommand=self.canvas_x.set, yscrollcommand=self.canvas_y.set)
        self.canvas.grid(row=0, column=0, sticky="nsew")
        self.canvas_y.grid(row=0, column=1, sticky="ns")
        self.canvas_x.grid(row=1, column=0, sticky="ew")
        canvas_wrap.rowconfigure(0, weight=1)
        canvas_wrap.columnconfigure(0, weight=1)

        # RIGHT CROP PANEL
        right = tk.Frame(content, bg="#f7fbff", width=560)
        right.pack(side="right", fill="both", padx=(6,0))

        tk.Label(right, text="Crop Field / Preview Before Export", bg="#f7fbff", font=("Arial",13,"bold")).pack(pady=3)

        rzoom = tk.Frame(right, bg="#f7fbff")
        rzoom.pack(fill="x")
        tk.Label(rzoom, text="Crop zoom 0–100%", bg="#f7fbff").pack(side="left")
        tk.Scale(rzoom, from_=0, to=100, orient="horizontal", variable=self.right_zoom, command=lambda e:self.refresh_right(), length=180).pack(side="left")
        tk.Button(rzoom, text="Fit", command=lambda:self.set_right_zoom(50)).pack(side="left", padx=3)
        tk.Button(rzoom, text="100%", command=lambda:self.set_right_zoom(100)).pack(side="left", padx=3)

        thumb_wrap = tk.Frame(right, bg="#f7fbff")
        thumb_wrap.pack(fill="x", padx=4)
        self.thumb_canvas = tk.Canvas(thumb_wrap, width=520, height=340, bg="white", highlightthickness=1, highlightbackground="#bbbbbb")
        self.thumb_scroll_y = tk.Scrollbar(thumb_wrap, orient="vertical", command=self.thumb_canvas.yview)
        self.thumb_scroll_x = tk.Scrollbar(thumb_wrap, orient="horizontal", command=self.thumb_canvas.xview)
        self.thumb_canvas.configure(yscrollcommand=self.thumb_scroll_y.set, xscrollcommand=self.thumb_scroll_x.set)
        self.thumb_canvas.grid(row=0, column=0, sticky="nsew")
        self.thumb_scroll_y.grid(row=0, column=1, sticky="ns")
        self.thumb_scroll_x.grid(row=1, column=0, sticky="ew")
        thumb_wrap.rowconfigure(0, weight=1)
        thumb_wrap.columnconfigure(0, weight=1)
        self.thumb_frame = tk.Frame(self.thumb_canvas, bg="white")
        self.thumb_canvas.create_window((0,0), window=self.thumb_frame, anchor="nw")
        self.thumb_frame.bind("<Configure>", lambda e:self.thumb_canvas.configure(scrollregion=self.thumb_canvas.bbox("all")))

        edit = tk.LabelFrame(right, text="Edit ACTIVE crop", bg="#f7fbff")
        edit.pack(fill="x", padx=5, pady=5)
        tk.Label(edit, text="Active ID", bg="#f7fbff").grid(row=0,column=0)
        self.id_var = tk.StringVar()
        tk.Entry(edit, textvariable=self.id_var, width=5).grid(row=0,column=1)
        self.btn(edit,"Set",None,self.set_active_from_entry).grid(row=0,column=2,padx=2)
        self.btn(edit,"Rotate -5°",None,lambda:self.rotate_active(-5)).grid(row=0,column=3,padx=2)
        self.btn(edit,"Rotate +5°",None,lambda:self.rotate_active(5)).grid(row=0,column=4,padx=2)
        self.btn(edit,"Reset",None,self.reset_active_angle).grid(row=0,column=5,padx=2)
        tk.Label(edit, text="Angle°", bg="#f7fbff").grid(row=0,column=6)
        tk.Entry(edit, textvariable=self.crop_angle_var, width=6).grid(row=0,column=7)
        self.btn(edit,"Apply angle",None,self.apply_active_angle).grid(row=0,column=8,padx=2)

        self.crop_erase_btn = self.btn(edit,"Erase crop brush", "#ffe1d9", self.set_erase_crop_mode)
        self.crop_erase_btn.grid(row=1,column=0,columnspan=2,sticky="ew",pady=2)
        self.crop_lasso_btn = self.btn(edit,"Erase crop circle/lasso", "#ffdfef", self.set_erase_crop_circle_mode)
        self.crop_lasso_btn.grid(row=1,column=2,columnspan=2,sticky="ew",pady=2)
        tk.Label(edit, text="Crop brush", bg="#f7fbff").grid(row=1,column=4)
        tk.Entry(edit, textvariable=self.crop_brush_size, width=4).grid(row=1,column=5)
        self.btn(edit,"Clear crop edit",None,self.clear_active_crop_edit).grid(row=2,column=0,columnspan=2,sticky="ew",pady=2)
        tk.Label(edit, text="Reassign to", bg="#f7fbff").grid(row=2,column=2)
        self.reassign_var = tk.StringVar()
        tk.Entry(edit, textvariable=self.reassign_var, width=5).grid(row=2,column=3)
        self.btn(edit,"Apply number",None,self.reassign_active_number).grid(row=2,column=4,columnspan=2,sticky="ew",pady=2)
        self.btn(edit,"Undo crop ↶", "#fff2cc", self.undo).grid(row=3,column=0,columnspan=2,sticky="ew",pady=2)
        self.btn(edit,"Redo crop ↷", "#fff2cc", self.redo).grid(row=3,column=2,columnspan=2,sticky="ew",pady=2)

        tk.Label(right, text="Active Crop Editor: Remove artifacts before export", bg="#f7fbff", font=("Arial",10,"bold")).pack()
        self.crop_canvas = tk.Canvas(right, width=550, height=350, bg="white", highlightthickness=1, highlightbackground="#bbbbbb")
        self.crop_canvas.pack(fill="both", expand=True, padx=5, pady=4)

        self.canvas.bind("<ButtonPress-1>", self.on_canvas_press)
        self.canvas.bind("<B1-Motion>", self.on_canvas_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_canvas_release)
        self.canvas.bind("<Button-3>", self.pick_active)

        self.crop_canvas.bind("<ButtonPress-1>", self.on_crop_press)
        self.crop_canvas.bind("<B1-Motion>", self.on_crop_drag)
        self.crop_canvas.bind("<ButtonRelease-1>", self.on_crop_release)

        self.root.bind("<Control-z>", lambda e:self.undo())
        self.root.bind("<Control-y>", lambda e:self.redo())
        self.root.bind("<Escape>", lambda e:self.set_toggle_mode())
        self.root.bind("<Left>", lambda e:self.rotate_active(-5))
        self.root.bind("<Right>", lambda e:self.rotate_active(5))

    def btn(self, parent, text, bg, cmd):
        b = tk.Button(parent, text=text, command=cmd)
        if bg: b.config(bg=bg)
        return b

    def set_left_zoom(self,val):
        self.left_zoom.set(val); self.refresh_canvas()
    def set_right_zoom(self,val):
        self.right_zoom.set(val); self.refresh_right()
    def zoom_factor(self,value):
        return 0.25 + (value/100)*1.75

    def snapshot(self):
        self.undo_stack.append((copy.deepcopy(self.candidates),
                                None if self.review_img is None else self.review_img.copy(),
                                copy.deepcopy(self.crop_edits), copy.deepcopy(self.crop_masks)))
        if len(self.undo_stack)>60: self.undo_stack.pop(0)
        self.redo_stack.clear()

    def undo(self):
        if not self.undo_stack:
            self.status.config(text="Nothing to undo."); return
        self.redo_stack.append((copy.deepcopy(self.candidates), None if self.review_img is None else self.review_img.copy(), copy.deepcopy(self.crop_edits), copy.deepcopy(self.crop_masks)))
        self.candidates, self.review_img, self.crop_edits, self.crop_masks = self.undo_stack.pop()
        self.refresh_all()
        self.status.config(text="Undo")

    def redo(self):
        if not self.redo_stack:
            self.status.config(text="Nothing to redo."); return
        self.undo_stack.append((copy.deepcopy(self.candidates), None if self.review_img is None else self.review_img.copy(), copy.deepcopy(self.crop_edits), copy.deepcopy(self.crop_masks)))
        self.candidates, self.review_img, self.crop_edits, self.crop_masks = self.redo_stack.pop()
        self.refresh_all()
        self.status.config(text="Redo")

    def clean_path(self,data):
        return data.strip().strip("{}").strip('"')
    def drop(self,event):
        self.image_path = Path(self.clean_path(event.data))
        self.dropbox.config(text=f"Loaded:\n{self.image_path.name}")
    def choose(self,event=None):
        path = filedialog.askopenfilename(filetypes=[("Images","*.jpg *.jpeg *.png *.tif *.tiff *.bmp")])
        if path:
            self.image_path = Path(path)
            self.dropbox.config(text=f"Loaded:\n{self.image_path.name}")

    def detect(self):
        if not self.image_path:
            messagebox.showwarning("No image","Please choose image first."); return
        self.original = imread_unicode(self.image_path)
        if self.original is None:
            messagebox.showerror("Error","Cannot open image."); return
        self.review_img = self.original.copy()
        self.candidates = detect_candidates(self.review_img, int(self.max_var.get()))
        self.crop_edits.clear(); self.crop_masks.clear(); self.crop_masks.clear(); self.undo_stack.clear(); self.redo_stack.clear()
        h,w = self.original.shape[:2]
        self.info_label.config(text=f"{w} x {h} pixels\nScale: 100% original\nNo resize applied")
        self.refresh_all()
        self.status.config(text=f"Detected {len(self.candidates)} candidates after filtering round stain/cell blobs. Use Keep/Reject, Add Box, then Auto Number.")

    def select_top_n(self):
        if not self.candidates: return
        self.snapshot()
        n = int(self.n_var.get())
        # Keep current order, select top N, then renumber selected as 1..N.
        for i,c in enumerate(self.candidates): 
            c["selected"] = i < n
        renumber(self.candidates)
        self.crop_edits.clear(); self.crop_masks.clear()
        self.refresh_all()
        self.status.config(text=f"Selected and synchronized numbers: original labels 1–{n} match crop field 1–{n}.")

    def set_mode_buttons(self, active):
        for b in [self.toggle_btn,self.add_btn,self.erase_btn,self.crop_erase_btn,self.crop_lasso_btn]:
            b.config(bg="SystemButtonFace")
        if active=="toggle": self.toggle_btn.config(bg="#d9ffd9")
        if active=="add": self.add_btn.config(bg="#d9e8ff")
        if active=="erase_original": self.erase_btn.config(bg="#ffe1d9")
        if active=="erase_crop": self.crop_erase_btn.config(bg="#ffe1d9")
        if active=="erase_crop_circle": self.crop_lasso_btn.config(bg="#ffdfef")

    def set_toggle_mode(self):
        self.mode.set("toggle"); self.set_mode_buttons("toggle"); self.status.config(text="Mode: Keep [Green] | Reject [Red] | Add Box Manually [Blue]  Click a chromosome box to change its status.")
    def set_add_mode(self):
        self.mode.set("add"); self.set_mode_buttons("add"); self.status.config(text="Mode: Add Box Manually [Blue]. Drag to add a missing chromosome box.")
    def set_erase_original_mode(self):
        self.mode.set("erase_original"); self.set_mode_buttons("erase_original"); self.status.config(text="Mode: Erase Original. Drag on the original image to remove unwanted marks.")
    def set_erase_crop_mode(self):
        self.mode.set("erase_crop"); self.set_mode_buttons("erase_crop"); self.status.config(text="Mode: Erase Crop Brush. Drag on the Active Crop Editor to remove artifacts.")
    def set_erase_crop_circle_mode(self):
        self.mode.set("erase_crop_circle"); self.set_mode_buttons("erase_crop_circle"); self.status.config(text="Mode: Erase Crop Circle/Lasso. Draw around the area to remove in the Active Crop Editor.")

    def refresh_all(self):
        self.refresh_canvas(); self.refresh_right()
    def refresh_right(self):
        self.refresh_thumbnails(); self.refresh_active_crop()

    def refresh_canvas(self):
        if self.review_img is None: return
        pil = draw_overlay(self.review_img, self.candidates, self.show_numbers.get())
        max_w,max_h = 850,720
        fit = min(max_w/pil.width, max_h/pil.height, 1.0)
        scale = fit * self.zoom_factor(self.left_zoom.get())
        self.display_scale = scale
        show = pil.resize((max(1,int(pil.width*scale)), max(1,int(pil.height*scale))))
        self.tk_img = ImageTk.PhotoImage(show)
        self.canvas.delete("all")
        self.canvas.create_image(0,0,image=self.tk_img,anchor="nw")
        self.canvas.configure(scrollregion=(0,0,show.width,show.height))

    def get_crop_for_candidate(self,c):
        # Always regenerate from current original + angle, then apply persistent erase mask.
        # This fixes: erase crop -> rotate -> erased part coming back.
        base_crop = crop_from_original(self.review_img, c, pad=40)

        if c["id"] in self.crop_masks:
            mask = self.crop_masks[c["id"]]
            if mask.shape[:2] != base_crop.shape[:2]:
                mask = cv2.resize(mask, (base_crop.shape[1], base_crop.shape[0]), interpolation=cv2.INTER_NEAREST)
                self.crop_masks[c["id"]] = mask
            base_crop[mask > 0] = (255,255,255)

        if c["id"] in self.crop_edits:
            edited = self.crop_edits[c["id"]]
            if edited.shape[:2] == base_crop.shape[:2]:
                return edited.copy()

        return base_crop

    def refresh_thumbnails(self):
        for child in self.thumb_frame.winfo_children(): child.destroy()
        self.thumb_refs.clear()
        if self.review_img is None: return
        z = self.zoom_factor(self.right_zoom.get())
        for idx,c in enumerate(selected_sorted(self.candidates)):
            crop = self.get_crop_for_candidate(c)
            pil = Image.fromarray(cv2.cvtColor(crop, cv2.COLOR_BGR2RGB))
            pil.thumbnail((max(1,int(120*z)), max(1,int(105*z))))
            tkimg = ImageTk.PhotoImage(pil); self.thumb_refs.append(tkimg)
            cell = tk.Frame(self.thumb_frame, bg="white", relief="ridge", bd=1)
            r,col = divmod(idx,4)
            cell.grid(row=r,column=col,padx=4,pady=4)
            lab = tk.Label(cell, text=f"Original #{c['id']}   {c.get('angle',0):.0f}°", bg="white", fg="green")
            lab.pack()
            img_lab = tk.Label(cell, image=tkimg, bg="white")
            img_lab.pack()
            img_lab.bind("<Button-1>", lambda e, cid=c["id"]: self.set_active_id(cid))
            lab.bind("<Button-1>", lambda e, cid=c["id"]: self.set_active_id(cid))

    def refresh_active_crop(self):
        self.crop_canvas.delete("all")
        c = self.get_active_candidate()
        if c is None or self.review_img is None: return
        crop = self.get_crop_for_candidate(c)
        pil = Image.fromarray(cv2.cvtColor(crop, cv2.COLOR_BGR2RGB))
        z = self.zoom_factor(self.right_zoom.get())
        fit = min(540/pil.width, 340/pil.height, 1.0)
        scale = fit*z
        self.crop_display_scale = scale
        show = pil.resize((max(1,int(pil.width*scale)), max(1,int(pil.height*scale))))
        self.crop_tk_img = ImageTk.PhotoImage(show)
        self.crop_canvas.create_image(5,5,image=self.crop_tk_img,anchor="nw")
        self.crop_canvas.create_text(12,12,text=f"#{c['id']}",fill="green",anchor="nw")

    def canvas_event_xy(self, x, y):
        # Convert visible-window mouse coordinates to full canvas coordinates.
        # This fixes click mismatch when the canvas is scrolled or zoomed.
        return self.canvas.canvasx(x), self.canvas.canvasy(y)

    def crop_canvas_event_xy(self, x, y):
        return self.crop_canvas.canvasx(x), self.crop_canvas.canvasy(y)

    def canvas_to_img(self,x,y):
        cx, cy = self.canvas_event_xy(x, y)
        return cx/self.display_scale, cy/self.display_scale

    def crop_canvas_to_img(self,x,y):
        cx, cy = self.crop_canvas_event_xy(x, y)
        return (cx-5)/self.crop_display_scale, (cy-5)/self.crop_display_scale

    def find_hit(self,x,y):
        margin=12
        hits=[c for c in reversed(self.candidates) if c["x"]-margin<=x<=c["x"]+c["w"]+margin and c["y"]-margin<=y<=c["y"]+c["h"]+margin]
        return hits[0] if hits else None

    def on_canvas_press(self,event):
        if self.review_img is None: return
        if self.mode.get()=="toggle":
            x,y=self.canvas_to_img(event.x,event.y)
            hit=self.find_hit(x,y)
            if hit:
                self.snapshot(); hit["selected"]=not hit["selected"]; renumber(self.candidates); self.refresh_all()
            return
        if self.mode.get()=="add":
            cx, cy = self.canvas_event_xy(event.x, event.y)
            self.drag_start=(cx,cy)
            self.temp_rect=self.canvas.create_rectangle(cx,cy,cx,cy,outline="blue",width=3)
            return
        if self.mode.get()=="erase_original":
            if not self.erase_snapshot_taken:
                self.snapshot(); self.erase_snapshot_taken=True
            self.erase_original_at(event.x,event.y)

    def on_canvas_drag(self,event):
        if self.mode.get()=="add" and self.drag_start and self.temp_rect:
            x0,y0=self.drag_start; cx, cy = self.canvas_event_xy(event.x, event.y); self.canvas.coords(self.temp_rect,x0,y0,cx,cy)
        elif self.mode.get()=="erase_original":
            self.erase_original_at(event.x,event.y)

    def on_canvas_release(self,event):
        if self.mode.get()=="erase_original":
            self.erase_snapshot_taken=False; self.crop_edits.clear(); self.crop_masks.clear(); self.refresh_all(); return
        if self.mode.get()!="add" or not self.drag_start: return
        x0,y0=self.drag_start; x1,y1=self.canvas_event_xy(event.x,event.y); self.drag_start=None
        if self.temp_rect: self.canvas.delete(self.temp_rect); self.temp_rect=None
        ix0,iy0=min(x0,x1)/self.display_scale, min(y0,y1)/self.display_scale
        ix1,iy1=max(x0,x1)/self.display_scale, max(y0,y1)/self.display_scale
        bw,bh=ix1-ix0, iy1-iy0
        if bw<10 or bh<10: return
        self.snapshot()
        new=dict(id=len(self.candidates)+1,x=int(ix0),y=int(iy0),w=int(bw),h=int(bh),cx=ix0+bw/2,cy=iy0+bh/2,area=int(bw*bh),length=int(max(bw,bh)),score=0,selected=True,manual=True,angle=0.0)
        self.candidates.append(new); renumber(self.candidates); self.refresh_all()

    def erase_original_at(self,x,y):
        ix,iy=self.canvas_to_img(x,y)
        r=int(self.brush_size.get())
        cv2.circle(self.review_img,(int(ix),int(iy)),r,(255,255,255),-1)
        self.refresh_canvas()

    def on_crop_press(self,event):
        if self.mode.get() not in ["erase_crop","erase_crop_circle"]: return
        if not self.erase_snapshot_taken:
            self.snapshot(); self.erase_snapshot_taken=True
        if self.mode.get()=="erase_crop":
            self.erase_crop_at(event.x,event.y)
        else:
            self.crop_lasso=[(event.x,event.y)]
            if self.crop_lasso_item: self.crop_canvas.delete(self.crop_lasso_item)
            self.crop_lasso_item=None

    def on_crop_drag(self,event):
        if self.mode.get()=="erase_crop":
            self.erase_crop_at(event.x,event.y)
        elif self.mode.get()=="erase_crop_circle":
            self.crop_lasso.append((event.x,event.y))
            if self.crop_lasso_item: self.crop_canvas.delete(self.crop_lasso_item)
            flat=[v for p in self.crop_lasso for v in p]
            if len(flat)>=4:
                self.crop_lasso_item=self.crop_canvas.create_line(*flat, fill="red", width=2, smooth=True)

    def on_crop_release(self,event):
        if self.mode.get()=="erase_crop_circle" and self.crop_lasso:
            self.apply_crop_lasso()
        self.erase_snapshot_taken=False
        self.crop_lasso=[]

    def erase_crop_at(self,x,y):
        c=self.get_active_candidate()
        if c is None: return
        crop=self.get_crop_for_candidate(c)
        ix,iy=self.crop_canvas_to_img(x,y)
        r=int(self.crop_brush_size.get())

        if c["id"] not in self.crop_masks or self.crop_masks[c["id"]].shape[:2] != crop.shape[:2]:
            self.crop_masks[c["id"]] = np.zeros(crop.shape[:2], dtype=np.uint8)

        cv2.circle(self.crop_masks[c["id"]], (int(ix),int(iy)), r, 255, -1)
        crop[self.crop_masks[c["id"]] > 0] = (255,255,255)
        self.crop_edits[c["id"]] = crop
        self.refresh_active_crop(); self.refresh_thumbnails()

    def apply_crop_lasso(self):
        c=self.get_active_candidate()
        if c is None or len(self.crop_lasso)<3: return
        crop=self.get_crop_for_candidate(c)
        pts=[]
        for x,y in self.crop_lasso:
            ix,iy=self.crop_canvas_to_img(x,y)
            pts.append([int(ix),int(iy)])
        pts=np.array([pts], dtype=np.int32)
        mask=np.zeros(crop.shape[:2], dtype=np.uint8)
        cv2.fillPoly(mask, pts, 255)
        if c["id"] not in self.crop_masks or self.crop_masks[c["id"]].shape[:2] != crop.shape[:2]:
            self.crop_masks[c["id"]] = np.zeros(crop.shape[:2], dtype=np.uint8)
        self.crop_masks[c["id"]][mask>0] = 255
        crop[self.crop_masks[c["id"]]>0]=(255,255,255)
        self.crop_edits[c["id"]]=crop
        if self.crop_lasso_item: self.crop_canvas.delete(self.crop_lasso_item); self.crop_lasso_item=None
        self.refresh_active_crop(); self.refresh_thumbnails()

    def pick_active(self,event):
        x,y=self.canvas_to_img(event.x,event.y)
        hit=self.find_hit(x,y)
        if hit: self.set_active_id(hit["id"])

    def set_active_id(self,cid):
        self.active_id=cid; self.id_var.set(str(cid))
        c = self.get_active_candidate()
        if c is not None:
            self.crop_angle_var.set(str(round(c.get("angle", 0), 2)))
        self.refresh_active_crop(); self.status.config(text=f"Active crop #{cid}")
    def set_active_from_entry(self):
        try: self.set_active_id(int(self.id_var.get()))
        except: messagebox.showwarning("No ID","Enter chromosome ID first.")
    def get_active_candidate(self):
        try: cid=int(self.id_var.get())
        except: cid=self.active_id
        for c in self.candidates:
            if c["id"]==cid and c["selected"]: return c
        return None

    def rotate_active(self,deg):
        c=self.get_active_candidate()
        if c is None: self.status.config(text="No active selected crop."); return
        self.snapshot()
        c["angle"] = c.get("angle",0) + deg
        self.crop_angle_var.set(str(round(c.get("angle",0), 2)))
        # Keep crop erase masks when rotating. The mask is resized to the new rotated crop,
        # so erased debris will not come back after rotation.
        self.crop_edits.pop(c["id"], None)
        self.refresh_all()

    def apply_active_angle(self):
        c=self.get_active_candidate()
        if c is None: self.status.config(text="No active selected crop."); return
        try:
            angle=float(self.crop_angle_var.get())
        except Exception:
            messagebox.showwarning("Invalid angle", "Please enter a number, e.g. -12.5 or 30")
            return
        self.snapshot()
        c["angle"] = angle
        self.crop_edits.pop(c["id"], None)
        self.refresh_all()

    def reset_active_angle(self):
        c=self.get_active_candidate()
        if c is None: return
        self.snapshot(); c["angle"]=0; self.crop_angle_var.set("0"); self.crop_edits.pop(c["id"],None); self.refresh_all()
    def clear_active_crop_edit(self):
        c=self.get_active_candidate()
        if c is None: return
        self.snapshot(); self.crop_edits.pop(c["id"],None); self.crop_masks.pop(c["id"],None); self.refresh_all()

    def delete_rejected(self):
        self.snapshot(); self.candidates=[c for c in self.candidates if c["selected"]]; renumber(self.candidates); self.crop_edits.clear(); self.crop_masks.clear(); self.refresh_all()
    def assign_length(self):
        self.snapshot(); assign_by_length(self.candidates); self.crop_edits.clear(); self.crop_masks.clear(); self.refresh_all()
        self.status.config(text="Auto Number (Chromosome Size): larger chromosomes are numbered first. Original labels match the Crop Field.")
    def assign_position(self):
        self.snapshot(); assign_by_position(self.candidates); self.crop_edits.clear(); self.crop_masks.clear(); self.refresh_all()
        self.status.config(text="Auto Number (Position): chromosomes are numbered from top-left to bottom-right. Original labels match the Crop Field.")

    def reassign_active_number(self):
        c=self.get_active_candidate()
        if c is None: return
        try: new_id=int(self.reassign_var.get())
        except: messagebox.showwarning("No number","Enter new number."); return
        selected=selected_sorted(self.candidates)
        old_idx=selected.index(c); new_idx=max(0,min(new_id-1,len(selected)-1))
        self.snapshot(); selected.pop(old_idx); selected.insert(new_idx,c)
        rejected=[x for x in self.candidates if not x["selected"]]
        self.candidates[:]=selected+rejected; renumber(self.candidates); self.crop_edits.clear(); self.crop_masks.clear(); self.set_active_id(new_idx+1); self.refresh_all()



    def open_help_window(self):
        """Open quick help and offer the full PDF manual."""
        win = tk.Toplevel(self.root)
        win.title("Cute Auto Karyotype [Quick Help]")
        win.geometry("760x680")
        win.configure(bg="white")
        txt = tk.Text(win, wrap="word", font=("Arial", 11), bg="white")
        txt.pack(fill="both", expand=True, padx=10, pady=10)
        help_text = (
            "Cute Auto Karyotype [Quick Help]\n\n"
            "Recommended Workflow\n"
            "1. Click Detect to automatically identify chromosomes in the image.\n"
            "2. Review each detected box and click to Keep or Reject it as needed. Use Add Box to manually add any missing chromosome regions.\n"
            "3. Click Delete Rejected to remove all rejected boxes from the workspace.\n"
            "4. Click Auto Number to automatically assign chromosome numbers based on the selected sorting criteria.\n"
            "5. Click a chromosome image in the Crop Field (upper-right panel). The selected chromosome will appear in the Active Crop Editor (lower-right panel).\n"
            "6. Remove unwanted artifacts or debris in the crop using the Brush or Lasso tools.\n"
            "7. Rotate the chromosome using the rotation buttons or enter a custom rotation angle manually.\n"
            "8. Export the results as PNG or PDF, or further arrange the chromosomes using Create Final Board to generate the final karyotype layout.\n"
            "9. Before exporting the final results, carefully review chromosome detection, artifact removal, orientation, and numbering to ensure the karyotype is accurate and properly organized.\n"
            "10. Note: The full manual is a separate PDF in the same folder.\n\n\n"
            "ขั้นตอนการทำงานที่แนะนำ\n"
            "1. กด Detect เพื่อให้โปรแกรมตรวจหาโครโมโซมในภาพโดยอัตโนมัติ\n"
            "2. ตรวจสอบกล่องที่ตรวจพบ แล้วคลิกเลือก Keep หรือ Reject ตามความเหมาะสม หากมีโครโมโซมที่ตรวจไม่พบ สามารถใช้ Add Box เพื่อเพิ่มกล่องด้วยตนเอง\n"
            "3. กด Delete Rejected เพื่อลบกล่องที่ถูกปฏิเสธออกจากพื้นที่ทำงาน\n"
            "4. กด Auto Number เพื่อจัดเรียงและกำหนดหมายเลขโครโมโซมโดยอัตโนมัติ\n"
            "5. คลิกภาพโครโมโซมใน Crop Field (แผงด้านขวาบน) ภาพที่เลือกจะแสดงใน Active Crop Editor (แผงด้านขวาล่าง)\n"
            "6. ลบเศษสิ่งรบกวนหรือส่วนที่ไม่ต้องการในภาพครอปด้วยเครื่องมือ Brush หรือ Lasso\n"
            "7. หมุนโครโมโซมด้วยปุ่มหมุน หรือกำหนดมุมการหมุนเองโดยการกรอกค่าองศา\n"
            "8. ส่งออกผลลัพธ์เป็นไฟล์ PNG หรือ PDF หรือจัดเรียงโครโมโซมเพิ่มเติมด้วย Create Final Board เพื่อสร้างภาพคาริโอไทป์ฉบับสมบูรณ์\n"
            "9. ควรตรวจสอบการตรวจจับ การลบสิ่งรบกวน และการจัดเรียงหมายเลขโครโมโซมให้ถูกต้องก่อนส่งออกผลลัพธ์ขั้นสุดท้าย\n"
            "10. หมายเหตุ: คู่มือการใช้งานฉบับสมบูรณ์อยู่ในไฟล์ PDF แยกต่างหากภายในโฟลเดอร์เดียวกันกับโปรแกรม\n"
        )
        txt.insert("1.0", help_text)
        txt.config(state="disabled")
        bar = tk.Frame(win, bg="white")
        bar.pack(fill="x", padx=10, pady=(0,10))
        tk.Button(bar, text="Open Full Manual PDF", bg="#d9f0ff", command=self.open_manual_pdf).pack(side="left", padx=4)
        tk.Button(bar, text="Close", command=win.destroy).pack(side="right", padx=4)

    def open_about_window(self):
        messagebox.showinfo(
        "About Cute Auto Karyotype",
        "Cute Auto Karyotype✨\n\n"
        "Version 1.0\n\n"
        "Developed by\n"
        "Praopilas Phakdeedindan\n\n"
        "Department of Animal Husbandry\n"
        "Faculty of Veterinary Science\n"
        "Chulalongkorn University\n"
        "Bangkok, Thailand\n\n"
        "Academic and Non-Commercial Use\n\n"
        "Copyright © 2026\n"
        "Praopilas Phakdeedindan\n\n"
        "All Rights Reserved."
        )

    def open_manual_pdf(self):
        manual = Path(__file__).with_name("Cute_Auto_Karyotype_Manual.pdf")
        if not manual.exists():
            messagebox.showwarning("Manual not found", f"Cannot find manual PDF:\n{manual}")
            return
        try:
            os.startfile(str(manual))
        except Exception as e:
            messagebox.showerror("Cannot open manual", str(e))

    # =======================
    # V17 FINAL BOARD EDITOR
    # =======================
    def crop_to_rgba(self, crop):
        """Convert BGR crop with white background to RGBA transparent layer."""
        rgb = cv2.cvtColor(crop, cv2.COLOR_BGR2RGB)
        pil = Image.fromarray(rgb).convert("RGBA")
        arr = np.array(pil)
        # near-white background becomes transparent
        white = (arr[:,:,0] > 246) & (arr[:,:,1] > 246) & (arr[:,:,2] > 246)
        arr[:,:,3] = np.where(white, 0, 255).astype(np.uint8)
        return Image.fromarray(arr)

    def auto_board_items(self):
        selected = selected_sorted(self.candidates)
        items = []
        if not selected:
            return items
        pairs_per_row = 4
        pair_w = 210
        row_h = 190
        start_x = 120
        start_y = 110
        for i, c in enumerate(selected):
            pair = i // 2
            within = i % 2
            row = pair // pairs_per_row
            col = pair % pairs_per_row
            x = start_x + col * pair_w + (within * 58)
            y = start_y + row * row_h
            crop = self.get_crop_for_candidate(c)
            rgba = self.crop_to_rgba(crop)
            # Keep every chromosome at the same original crop scale.
            # Do not normalize each piece individually; otherwise small chromosomes can look bigger than large ones.
            scale = 1.0
            items.append({"id": c["id"], "rgba": rgba, "x": x, "y": y, "scale": scale, "angle": 0.0})
        return items

    def create_final_board(self):
        if self.review_img is None or not selected_sorted(self.candidates):
            messagebox.showwarning("No data", "Detect and select chromosomes first.")
            return
        self.board_items = self.auto_board_items()
        if self.board_win is not None and self.board_win.winfo_exists():
            self.board_win.lift()
            self.refresh_board()
            return
        self.board_win = tk.Toplevel(self.root)
        self.board_win.title("Cute Auto Karyotype - Final Board")
        self.board_win.geometry("1150x820")
        top = tk.Frame(self.board_win, bg="#f8fbff")
        top.pack(fill="x")
        tk.Label(top, text="Final Board: fixed original scale; drag chromosomes, click to select, rotate, then Save", bg="#f8fbff", font=("Arial", 12, "bold")).pack(side="left", padx=8, pady=6)
        tk.Button(top, text="Auto arrange", command=self.board_auto_arrange).pack(side="left", padx=3)
        tk.Button(top, text="Rotate -5°", command=lambda:self.board_rotate(-5)).pack(side="left", padx=3)
        tk.Button(top, text="Rotate +5°", command=lambda:self.board_rotate(5)).pack(side="left", padx=3)
        tk.Label(top, text="Angle°", bg="#f8fbff").pack(side="left", padx=(8,2))
        self.board_angle_var = tk.StringVar(value="0")
        tk.Entry(top, textvariable=self.board_angle_var, width=6).pack(side="left", padx=2)
        tk.Button(top, text="Apply", command=self.board_apply_angle).pack(side="left", padx=3)
        tk.Button(top, text="Save final board", bg="#c9f7ff", command=self.save_final_board).pack(side="left", padx=8)
        wrap = tk.Frame(self.board_win)
        wrap.pack(fill="both", expand=True)
        self.board_canvas = tk.Canvas(wrap, width=1050, height=720, bg="white", scrollregion=(0,0,1500,1000))
        sx = tk.Scrollbar(wrap, orient="horizontal", command=self.board_canvas.xview)
        sy = tk.Scrollbar(wrap, orient="vertical", command=self.board_canvas.yview)
        self.board_canvas.configure(xscrollcommand=sx.set, yscrollcommand=sy.set)
        self.board_canvas.grid(row=0, column=0, sticky="nsew")
        sy.grid(row=0, column=1, sticky="ns")
        sx.grid(row=1, column=0, sticky="ew")
        wrap.rowconfigure(0, weight=1); wrap.columnconfigure(0, weight=1)
        self.board_canvas.bind("<ButtonPress-1>", self.board_press)
        self.board_canvas.bind("<B1-Motion>", self.board_drag_move)
        self.board_canvas.bind("<ButtonRelease-1>", self.board_release)
        self.refresh_board()

    def board_auto_arrange(self):
        self.board_items = self.auto_board_items()
        self.refresh_board()

    def board_render_item_image(self, item):
        img = item["rgba"]
        if item.get("angle", 0):
            img = img.rotate(item["angle"], expand=True, resample=Image.Resampling.BICUBIC)
        scale = item.get("scale", 1.0)
        if scale != 1.0:
            img = img.resize((max(1, int(img.width*scale)), max(1, int(img.height*scale))), Image.Resampling.LANCZOS)
        return img

    def refresh_board(self):
        if self.board_canvas is None:
            return
        self.board_canvas.delete("all")
        self.board_refs.clear()
        self.board_canvas.create_text(50, 30, text=f"{self.image_path.stem if self.image_path else 'Karyotype'}", anchor="nw", font=("Arial", 22, "bold"), fill="#111")
        self.board_canvas.create_text(50, 65, text=f"Selected chromosomes: {len(self.board_items)}", anchor="nw", font=("Arial", 12), fill="#333")
        # pair guide labels
        for pair in range(math.ceil(len(self.board_items)/2)):
            row = pair // 4; col = pair % 4
            px = 120 + col*210 + 28; py = 250 + row*190
            self.board_canvas.create_text(px, py, text=str(pair+1), font=("Arial", 13, "bold"), fill="#555")
        for item in self.board_items:
            img = self.board_render_item_image(item)
            tkimg = ImageTk.PhotoImage(img)
            self.board_refs.append(tkimg)
            tag = f"chr_{item['id']}"
            self.board_canvas.create_image(item["x"], item["y"], image=tkimg, anchor="center", tags=(tag, "chr"))
            color = "red" if self.board_active == item["id"] else "green"
            self.board_canvas.create_text(item["x"], item["y"]-70, text=f"#{item['id']}", fill=color, font=("Arial", 11, "bold"), tags=(tag, "chr"))

    def board_find_item(self, x, y):
        cx = self.board_canvas.canvasx(x); cy = self.board_canvas.canvasy(y)
        best = None; best_d = 10**9
        for item in self.board_items:
            d = (item["x"]-cx)**2 + (item["y"]-cy)**2
            if d < best_d and d < 120**2:
                best = item; best_d = d
        return best, cx, cy

    def board_press(self, event):
        item, cx, cy = self.board_find_item(event.x, event.y)
        if item is None:
            self.board_active = None; self.refresh_board(); return
        self.board_active = item["id"]
        if hasattr(self, "board_angle_var"):
            self.board_angle_var.set(str(round(item.get("angle",0), 2)))
        self.board_drag = (item["id"], cx - item["x"], cy - item["y"])
        self.refresh_board()

    def board_drag_move(self, event):
        if not self.board_drag:
            return
        cid, dx, dy = self.board_drag
        cx = self.board_canvas.canvasx(event.x); cy = self.board_canvas.canvasy(event.y)
        for item in self.board_items:
            if item["id"] == cid:
                item["x"] = cx - dx; item["y"] = cy - dy
                break
        self.refresh_board()

    def board_release(self, event):
        self.board_drag = None

    def board_rotate(self, deg):
        if self.board_active is None:
            return
        for item in self.board_items:
            if item["id"] == self.board_active:
                item["angle"] = item.get("angle", 0) + deg
                if hasattr(self, "board_angle_var"):
                    self.board_angle_var.set(str(round(item.get("angle", 0), 2)))
                break
        self.refresh_board()

    def board_apply_angle(self):
        if self.board_active is None:
            return
        try:
            angle = float(self.board_angle_var.get())
        except Exception:
            messagebox.showwarning("Invalid angle", "Please enter a number, e.g. -12.5 or 30")
            return
        for item in self.board_items:
            if item["id"] == self.board_active:
                item["angle"] = angle
                break
        self.refresh_board()

    def board_scale(self, factor):
        if self.board_active is None:
            return
        for item in self.board_items:
            if item["id"] == self.board_active:
                item["scale"] = max(0.15, min(3.0, item.get("scale", 1.0) * factor))
                break
        self.refresh_board()

    def compose_final_board_image(self):
        W, H = 1500, 1000
        board = Image.new("RGBA", (W, H), (255,255,255,255))
        draw = ImageDraw.Draw(board)
        title = self.image_path.stem if self.image_path else "Karyotype"
        draw.text((50, 30), title, fill=(0,0,0,255), font=get_label_font(38))
        draw.text((50, 80), f"Selected chromosomes: {len(self.board_items)}", fill=(40,40,40,255), font=get_label_font(20))
        for pair in range(math.ceil(len(self.board_items)/2)):
            row = pair // 4; col = pair % 4
            draw.text((120 + col*210 + 22, 245 + row*190), str(pair+1), fill=(70,70,70,255), font=get_label_font(22))
        for item in self.board_items:
            img = self.board_render_item_image(item)
            x = int(item["x"] - img.width/2); y = int(item["y"] - img.height/2)
            board.alpha_composite(img, (x,y))
        return board.convert("RGB")

    def save_final_board(self):
        if not self.board_items:
            messagebox.showwarning("No board", "Create final board first.")
            return
        self.output_dir.mkdir(exist_ok=True)
        layers_dir = self.output_dir / "final_board_layers_transparent"
        layers_dir.mkdir(exist_ok=True)
        board = self.compose_final_board_image()
        png = self.output_dir / "final_karyotype_board.png"
        pdf = self.output_dir / "final_karyotype_board.pdf"
        board.save(png)
        board.save(pdf, "PDF", resolution=300.0)
        # Save each chromosome as a transparent layer with placement metadata.
        with open(self.output_dir / "final_board_layers.csv", "w", newline="", encoding="utf-8-sig") as f:
            wr = csv.writer(f); wr.writerow(["id", "x", "y", "scale", "angle", "file"])
            for item in self.board_items:
                img = self.board_render_item_image(item)
                layer_name = f"chr_{item['id']:02d}_transparent.png"
                img.save(layers_dir / layer_name)
                wr.writerow([item["id"], round(item["x"],2), round(item["y"],2), round(item["scale"],4), round(item["angle"],2), layer_name])
        # Optional PPTX output, each chromosome is a movable image object.
        try:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(5), Inches(0.5)).text = self.image_path.stem if self.image_path else "Karyotype"
            for item in self.board_items:
                layer_name = layers_dir / f"chr_{item['id']:02d}_transparent.png"
                # map 1500x1000 board to 13.33x7.5 slide approximately
                left = Inches(item["x"] / 1500 * 13.333)
                top = Inches(item["y"] / 1000 * 7.5)
                width = Inches(max(0.15, item["rgba"].width * item["scale"] / 1500 * 13.333))
                pic = slide.shapes.add_picture(str(layer_name), left, top, width=width)
                pic.left = int(left - pic.width/2)
                pic.top = int(top - pic.height/2)
            prs.save(self.output_dir / "final_karyotype_board_editable_layers.pptx")
        except Exception:
            pass
        messagebox.showinfo("Done", f"Saved final board and transparent layers\n{self.output_dir}")

    def export_all(self):
        if self.review_img is None or not self.candidates:
            messagebox.showwarning("No data","Detect first."); return
        self.output_dir.mkdir(exist_ok=True)
        crops_dir=self.output_dir/"chromosome_crops"; crops_dir.mkdir(exist_ok=True)
        layers_dir=self.output_dir/"chromosome_layers_transparent"; layers_dir.mkdir(exist_ok=True)
        for old in list(crops_dir.glob("chromosome_*.png")) + list(layers_dir.glob("chromosome_*.png")) :
            try: old.unlink()
            except: pass
        selected=selected_sorted(self.candidates)
        crop_items=[]
        for c in selected:
            crop=self.get_crop_for_candidate(c)
            crop_items.append((c["id"],crop))
            Image.fromarray(cv2.cvtColor(crop, cv2.COLOR_BGR2RGB)).save(crops_dir/f"chromosome_{c['id']:02d}.png")
            # V17 transparent crop layer
            trans_dir = self.output_dir/"chromosome_layers_transparent"; trans_dir.mkdir(exist_ok=True)
            self.crop_to_rgba(crop).save(trans_dir/f"chromosome_{c['id']:02d}_transparent.png")
            crop_to_transparent_rgba(crop).save(layers_dir/f"chromosome_{c['id']:02d}_transparent.png")

        max_w=max([crop.shape[1] for _,crop in crop_items]+[120])
        max_h=max([crop.shape[0] for _,crop in crop_items]+[160])
        cell_w=max_w+60; cell_h=max_h+70
        cols=min(7,max(1,len(crop_items))); rows=math.ceil(len(crop_items)/cols)
        sheet=Image.new("RGB",(cols*cell_w, rows*cell_h+95),"white")
        draw=ImageDraw.Draw(sheet)
        draw.text((30,30),f"Cute Auto Karyotype V17.2 | original-scale crops | selected: {len(crop_items)}",fill=(20,20,20))
        for i,(cid,crop) in enumerate(crop_items):
            r,col=divmod(i,cols); x0,y0=col*cell_w,85+r*cell_h
            draw.rounded_rectangle([x0+10,y0+10,x0+cell_w-10,y0+cell_h-10], radius=18, outline=(220,220,220), width=2, fill=(252,252,252))
            pil=Image.fromarray(cv2.cvtColor(crop, cv2.COLOR_BGR2RGB))
            sheet.paste(pil,(x0+(cell_w-pil.width)//2,y0+25))
            draw.text((x0+cell_w//2-6,y0+cell_h-35),str(cid),fill=(30,30,30))
        png=self.output_dir/"karyotype_sheet.png"; pdf=self.output_dir/"karyotype_sheet.pdf"; svg=self.output_dir/"karyotype_sheet_editable_layers.svg"; overlay=self.output_dir/"numbered_original_overlay.png"
        sheet.save(png); sheet.save(pdf,"PDF",resolution=300.0)
        make_svg_layer_file(svg, crop_items, cell_w, cell_h, cols)
        draw_overlay(self.review_img,self.candidates,True).save(overlay)
        with open(self.output_dir/"review_table.csv","w",newline="",encoding="utf-8-sig") as f:
            wr=csv.writer(f); wr.writerow(["id","selected","manual","x","y","w","h","angle","crop_edited"])
            for c in self.candidates:
                wr.writerow([c["id"],c["selected"],c.get("manual",False),c["x"],c["y"],c["w"],c["h"],c.get("angle",0),(c["id"] in self.crop_edits) or (c["id"] in self.crop_masks)])
        messagebox.showinfo("Done",f"Exported PNG/PDF/SVG + transparent layer PNGs\n{self.output_dir}")

    def open_folder(self):
        self.output_dir.mkdir(exist_ok=True); os.startfile(str(self.output_dir))

if __name__=="__main__":
    try:
        root=TkinterDnD.Tk() if DND_OK else tk.Tk()
        app=App(root)
        root.mainloop()
    except Exception as e:
        import traceback
        log_path = Path.cwd() / "v17_error_log.txt"
        log_path.write_text(traceback.format_exc(), encoding="utf-8")
        try:
            messagebox.showerror("Cute Auto Karyotype V17 error", f"Program could not start. Error saved to:\n{log_path}\n\n{e}")
        except Exception:
            print(traceback.format_exc())
        raise
